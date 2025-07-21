#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Xsentiment project
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_xsentiment_presentation():
    """Create a comprehensive PowerPoint presentation for Xsentiment"""
    
    # Create presentation object
    prs = Presentation()
    
    # Define color scheme
    primary_color = RGBColor(29, 161, 242)  # Twitter blue
    secondary_color = RGBColor(47, 79, 79)  # Dark slate gray
    accent_color = RGBColor(255, 107, 107)  # Light red
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "🐦 Xsentiment"
    subtitle.text = "Twitter Sentiment Analysis Tool\n\nReal-time Tweet Collection & Analysis\nComprehensive Visualization Dashboard"
    
    # Format title
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.color.rgb = primary_color
    
    # Format subtitle
    subtitle_para = subtitle.text_frame.paragraphs[0]
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = secondary_color
    
    # Slide 2: Project Overview
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "📊 Project Overview"
    content.text = """Purpose: Comprehensive Twitter sentiment analysis solution

Key Objectives:
• Monitor public opinion in real-time
• Classify tweet sentiment automatically
• Generate actionable insights through visualizations
• Support research and business intelligence

Target Audience:
• Social media analysts
• Market researchers
• Academic researchers
• Business intelligence teams"""
    
    # Slide 3: System Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "🏗️ System Architecture"
    content.text = """Data Flow Pipeline:

1. Data Collection Layer
   • Twitter API v2 integration
   • Real-time tweet streaming
   • Keyword-based filtering

2. Processing Layer
   • TextBlob sentiment analysis
   • Duplicate detection & removal
   • Data cleaning & preprocessing

3. Storage Layer
   • CSV-based data persistence
   • Timestamp tracking
   • Incremental data updates

4. Visualization Layer
   • Multiple chart types
   • Interactive dashboards
   • Export capabilities"""
    
    # Slide 4: Technical Features
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "⚡ Technical Features"
    content.text = """Core Capabilities:

🔄 Continuous Operation
• Runs indefinitely until interrupted
• 5-minute collection cycles
• Automatic rate limit handling

🧠 Sentiment Analysis
• TextBlob polarity scoring
• Three-class classification (Positive/Negative/Neutral)
• English language filtering

📊 Advanced Visualizations
• Static charts (PNG): Pie, Bar, Timeline
• Word clouds by sentiment
• Interactive Plotly dashboards
• Real-time data updates

🛡️ Robust Error Handling
• API rate limit management
• Graceful shutdown on interruption
• Comprehensive logging"""
    
    # Slide 5: Key Components
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "🔧 Key Components"
    content.text = """Main Scripts:

📄 script.py (Data Collection Engine)
• Twitter API integration with Tweepy
• Continuous tweet collection
• Real-time sentiment analysis
• CSV data persistence
• Duplicate prevention

📄 visualize_tweets.py (Visualization Engine)
• Multiple chart types generation
• Interactive dashboard creation
• Word cloud generation
• Timeline analysis
• Statistical summaries

📄 requirements_viz.txt
• All necessary Python dependencies
• Visualization libraries (Matplotlib, Plotly, Seaborn)
• Data processing tools (Pandas, NumPy)"""
    
    # Slide 6: Methodology
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "🔬 Methodology"
    content.text = """Step-by-Step Process:

1️⃣ Data Collection
   • Connect to Twitter API v2
   • Search recent tweets by keyword
   • Filter English-language tweets only
   • Handle pagination for large datasets

2️⃣ Sentiment Analysis
   • Apply TextBlob sentiment polarity
   • Convert polarity to categorical labels:
     - Positive: polarity > 0
     - Negative: polarity < 0
     - Neutral: polarity = 0

3️⃣ Data Storage
   • Append new tweets to CSV
   • Remove duplicates based on text
   • Add collection timestamps
   • Maintain data integrity

4️⃣ Visualization & Analysis
   • Generate multiple chart types
   • Create interactive dashboards
   • Provide statistical summaries"""
    
    # Slide 7: Usage Instructions
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "🚀 How to Use"
    content.text = """Setup & Configuration:

1. Prerequisites
   • Python 3.9+
   • Twitter Developer Account
   • API Bearer Token

2. Installation
   pip3 install tweepy textblob pandas
   pip3 install -r requirements_viz.txt

3. Configuration
   • Update BEARER_TOKEN in script.py
   • Modify KEYWORD for target search terms
   • Adjust MAX_RESULTS and SLEEP_INTERVAL

4. Execution
   # Start data collection
   python3 script.py
   
   # Generate visualizations
   python3 visualize_tweets.py

5. Stop Collection
   • Press Ctrl+C for graceful shutdown"""
    
    # Slide 8: Visualization Types
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "📈 Visualization Types"
    content.text = """Available Chart Types:

🥧 Pie Charts
• Sentiment distribution percentages
• Color-coded by sentiment type
• Clear percentage labels

📊 Bar Charts
• Tweet counts by sentiment
• Value labels on bars
• Professional styling

📅 Timeline Charts
• Sentiment trends over time
• Multi-line plots for each sentiment
• Date-based analysis

☁️ Word Clouds
• Most common words by sentiment
• Filtered stop words
• Visual word frequency representation

🎛️ Interactive Dashboard
• Multi-panel Plotly interface
• Hover effects and zoom
• Tweet length distribution
• Top words analysis"""
    
    # Slide 9: Configuration Options
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "⚙️ Configuration Options"
    content.text = """Customizable Parameters:

🔍 Search Configuration
KEYWORD = "OpenAI"           # Target search terms
MAX_RESULTS = 10             # Tweets per cycle
SLEEP_INTERVAL = 300         # Seconds between cycles

📝 Advanced Search Examples
• Single term: "ChatGPT"
• Multiple terms: "AI OR ML OR DeepLearning"
• Hashtags: "#MachineLearning"
• Complex queries: "OpenAI AND (GPT OR ChatGPT)"

⏰ Rate Limit Management
• Twitter API v2: 300 requests/15 minutes
• Built-in wait_on_rate_limit handling
• Configurable sleep intervals
• Automatic retry mechanism

💾 Data Management
• CSV file location: tweets_sentiment.csv
• Automatic backup and duplicate removal
• Timestamped data collection
• Long-term data persistence"""
    
    # Slide 10: Sample Output
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "📋 Sample Output"
    content.text = """Example Dataset Summary:

==================================================
DATASET SUMMARY
==================================================
Total tweets: 150
Columns: ['Text', 'Sentiment', 'Created_At', 'Collection_Time']

Sentiment Distribution:
  Positive: 89 (59.3%)
  Neutral: 45 (30.0%)
  Negative: 16 (10.7%)

Date range: 2025-06-14 to 2025-06-15

Sample tweets:
  [Positive] Amazing progress in AI development...
  [Neutral] New AI paper released today...
  [Negative] Concerns about AI safety..."""
    
    # Slide 11: Use Cases & Applications
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "🎯 Use Cases & Applications"
    content.text = """Real-World Applications:

🏢 Business Intelligence
• Brand sentiment monitoring
• Product launch feedback analysis
• Competitor sentiment tracking
• Crisis management response

🎓 Academic Research
• Public opinion studies
• Social media behavior analysis
• Political sentiment research
• Event impact assessment

📈 Marketing & PR
• Campaign effectiveness measurement
• Influencer impact analysis
• Trend identification
• Audience sentiment profiling

🔍 Market Research
• Consumer opinion analysis
• Product feedback collection
• Market trend identification
• Customer satisfaction monitoring"""
    
    # Slide 12: Technical Specifications
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "🔧 Technical Specifications"
    content.text = """System Requirements:

💻 Environment
• Python 3.9+
• macOS/Linux/Windows compatible
• Internet connection required

📚 Dependencies
Core Libraries:
• tweepy (Twitter API)
• textblob (Sentiment analysis)
• pandas (Data manipulation)

Visualization Libraries:
• matplotlib (Static plots)
• seaborn (Statistical plots)
• plotly (Interactive dashboards)
• wordcloud (Word cloud generation)

🔐 API Requirements
• Twitter Developer Account
• Bearer Token for API v2 access
• Rate limits: 300 requests/15 minutes

💾 Storage
• CSV format for data persistence
• Automatic file management
• Incremental data updates"""
    
    # Slide 13: Security & Best Practices
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "🔒 Security & Best Practices"
    content.text = """Security Considerations:

🔑 API Key Management
• Never commit API keys to version control
• Use environment variables for credentials
• Implement secure key rotation practices

💡 Recommended Setup
import os
BEARER_TOKEN = os.getenv('TWITTER_BEARER_TOKEN')

📁 Data Protection
• .gitignore includes sensitive data files
• Local CSV storage only
• No cloud upload by default

⚡ Performance Optimization
• Respect API rate limits
• Implement efficient caching
• Use appropriate sleep intervals
• Monitor resource usage

🛡️ Error Handling
• Graceful shutdown on interrupts
• Comprehensive logging
• Automatic retry mechanisms
• Data integrity validation"""
    
    # Slide 14: Future Enhancements
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "🔮 Future Enhancements"
    content.text = """Potential Improvements:

🤖 Advanced Analytics
• Machine learning-based sentiment analysis
• Emotion detection (joy, anger, fear, etc.)
• Sarcasm and context awareness
• Multi-language support

☁️ Cloud Integration
• Real-time streaming with Apache Kafka
• Cloud database storage (PostgreSQL, MongoDB)
• Scalable processing with Apache Spark
• Docker containerization

🌐 Web Dashboard
• Real-time web interface
• User authentication system
• Multi-user support
• Custom alert notifications

📱 Mobile Application
• iOS/Android companion apps
• Push notifications for trends
• Mobile-friendly visualizations
• Offline data viewing

🔧 Additional Features
• Influencer identification
• Hashtag trend analysis
• Geographic sentiment mapping
• Competitive analysis tools"""
    
    # Slide 15: Benefits & Impact
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "💡 Benefits & Impact"
    content.text = """Project Value Proposition:

✅ Cost-Effective Solution
• Open-source implementation
• No licensing fees
• Scalable architecture
• Minimal infrastructure requirements

📊 Data-Driven Insights
• Real-time sentiment tracking
• Historical trend analysis
• Quantitative opinion measurement
• Evidence-based decision making

🚀 Easy Implementation
• Simple setup process
• Comprehensive documentation
• Modular design
• Extensible architecture

🎯 Immediate Business Value
• Brand reputation monitoring
• Customer feedback analysis
• Market trend identification
• Crisis detection and response

🔬 Research Applications
• Academic research support
• Social media behavior studies
• Public opinion analysis
• Policy impact assessment"""
    
    # Slide 16: Demo & Screenshots
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "📸 Generated Visualizations"
    content.text = """Sample Output Files:

📊 Static Charts (PNG Format)
• sentiment_pie_chart.png
  - Circular representation of sentiment distribution
  - Percentage labels for each category
  - Color-coded segments

• sentiment_bar_chart.png
  - Bar chart with tweet counts
  - Value labels on each bar
  - Professional styling

• sentiment_timeline.png
  - Time-series analysis
  - Multi-line plot for trends
  - Date-based insights

☁️ Word Clouds
• wordcloud_positive.png
• wordcloud_negative.png
• wordcloud_neutral.png

🌐 Interactive Dashboard
• tweet_dashboard.html
  - Multi-panel interface
  - Interactive hover effects
  - Zoomable charts"""
    
    # Slide 17: Conclusion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "🎉 Conclusion"
    content.text = """Project Summary:

✨ What We Accomplished
• Built a comprehensive Twitter sentiment analysis tool
• Implemented real-time data collection and processing
• Created multiple visualization options
• Designed a user-friendly, extensible system

🎯 Key Strengths
• Real-time sentiment monitoring capabilities
• Professional visualization outputs
• Robust error handling and rate limit management
• Comprehensive documentation and setup instructions

🌟 Value Delivered
• Actionable insights from social media data
• Cost-effective alternative to commercial tools
• Educational resource for sentiment analysis
• Foundation for advanced analytics projects

🚀 Ready for Production
• Thoroughly tested components
• Scalable architecture
• Professional-grade visualizations
• Enterprise-ready documentation"""
    
    # Slide 18: Questions & Contact
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "❓ Questions & Discussion"
    content.text = """Thank you for your attention!

🤝 Let's Discuss
• Implementation questions
• Technical challenges
• Use case scenarios
• Future enhancement ideas

📧 Get Started
1. Clone the repository
2. Install dependencies
3. Configure Twitter API
4. Run your first analysis

🔗 Resources
• GitHub repository with full source code
• Comprehensive README documentation
• Installation and setup guides
• Example configurations

💬 Open Floor
• Questions about implementation?
• Interested in specific use cases?
• Want to discuss enhancements?
• Need help with setup?

Thank you! 🐦📊✨"""
    
    # Save presentation
    prs.save('Xsentiment_Comprehensive_Presentation.pptx')
    print("✅ PowerPoint presentation created successfully!")
    print("📄 File saved as: Xsentiment_Comprehensive_Presentation.pptx")
    print("\n📊 Presentation contains 18 comprehensive slides:")
    print("   1. Title Slide")
    print("   2. Project Overview")
    print("   3. System Architecture")
    print("   4. Technical Features")
    print("   5. Key Components")
    print("   6. Methodology")
    print("   7. Usage Instructions")
    print("   8. Visualization Types")
    print("   9. Configuration Options")
    print("   10. Sample Output")
    print("   11. Use Cases & Applications")
    print("   12. Technical Specifications")
    print("   13. Security & Best Practices")
    print("   14. Future Enhancements")
    print("   15. Benefits & Impact")
    print("   16. Demo & Screenshots")
    print("   17. Conclusion")
    print("   18. Questions & Contact")

if __name__ == "__main__":
    try:
        create_xsentiment_presentation()
    except ImportError:
        print("❌ python-pptx not installed!")
        print("📦 Please install it with: pip install python-pptx")
        print("🔄 Then run this script again.")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        print("🛠️ Please check the error and try again.")
